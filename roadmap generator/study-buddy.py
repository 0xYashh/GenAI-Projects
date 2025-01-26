import os 
import google.generativeai as genai
from dotenv import load_dotenv
from rich.console import Console
from rich.markdown import Markdown
from rich.progress import Progress
from PyPDF2 import PdfReader
from pptx import Presentation
from fpdf import FPDF

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
console = Console()

def clean_response(text):
    return text.replace("**","").replace("*","").strip()

def generate_roadmap(user_prompt, doc_content=""):
    system_prompt = f"""
    You are an expert study planner. Create a daily roadmap based on: {user_prompt}
    Document Context: {doc_content}
    **Strict Rules:**
    1. Format: Plain Text , NO MARKDOWN
    2. Structure: Day X (Total Hours): Topic: Time (Theory/Practice/Revision) Activity: Time (Exercise/Project)
    3. Content: Allocate ALL hours specified per day. Split time: Theory (40%), Practice (40%), Revision (20%).
    4. Example Output: Day 1 (5 hours): 1.thon Basics: 2 hours (Theory) 2.riables Practice: 2 hours (Practice) 3.de Exercises: 1 hour (Exercise) 4.oject Idea : "Build a calculator" 5.Project Idea : " Build a house predicting linear regression model"
    5. Resources: Add a "Resources" section at the end with: YouTube channels (e.g., "StatQuest for ML concepts")[Links]. Books (e.g., "Hands-On Machine Learning by Aurélien Géron")[Links]. Tools (e.g., "Jupyter Notebook, Scikit-learn")[Links]. Ensure resources are relevant and beginner-friendly.
    6. Creativity: Add engaging activities (e.g., "Build a spam email classifier"). Include practical projects (e.g., "Predict house prices using regression"). Suggest real-world datasets (e.g., "Titanic dataset for classification").
    """
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(system_prompt)
    return clean_response(response.text)

def process_document(file_path):
    try:
        if file_path.endswith('.pdf'):
            with open(file_path,'rb') as file:
                reader = PdfReader(file)
                doc_content = ' '.join([page.extract_text() for page in reader.pages]) 
                return doc_content
        elif file_path.endswith('.pptx'):
            prs = Presentation(file_path)
            doc_content = ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape,"text")]) 
            return doc_content  
    except Exception as e:
        console.print(f"[bold red]Error processing document: {e}[/]")
        return " "

def extract_tasks_from_roadmap(roadmap):
    """
    Extracts tasks from the roadmap based on the format provided.
    """
    tasks = []
    for line in roadmap.splitlines():
        if line.strip().startswith("Day"):
            tasks.append(line.strip())
        elif any(keyword in line.lower() for keyword in ["topic:", "activity:"]):
            tasks.append(line.strip())
    return '\n'.join(tasks)

def study_chatbot(roadmap="", doc_content=""):
    if not roadmap and not doc_content:
        console.print("[red]Generate a roadmap or upload a document first![/red]")
        return

    chat_session = genai.GenerativeModel('gemini-pro').start_chat(history=[])
    console.print("\n[bold green]Study Buddy: ASK ME ANYTHING ABOUT YOUR DOCUMENT AND ROADMAP![/] type 'exit' to quit\n")

    while True:
        user_input = console.input("[bold green]You: [/]").strip()
        if user_input.lower() == 'exit':
            break

        context = f"""
        Roadmap: {roadmap[:3000]}
        Document Content: {doc_content[:5000]}
        """
        prompt = f"""
        Answer based on this context: {context}
        
        Question: {user_input}
        """
        response = chat_session.send_message(prompt)
        console.print(Markdown(f"**Study Buddy:** {response.text}"))

def track_progress(roadmap):
    """
    Tracks user progress using a dynamically generated task list
    """
    # Extract tasks from roadmap
    tasks_text = extract_tasks_from_roadmap(roadmap)
    tasks = [line for line in tasks_text.split('\n') if line.strip()]
    completed = set()  # Use a set to avoid duplicates

    with Progress() as progress:
        task_progress = progress.add_task("[cyan]Progress:", total=len(tasks))

        while True:
            console.print("\n[bold green]Your Tasks:[/]")
            for i, task in enumerate(tasks, 1):
                status = "✅" if i-1 in completed else " "
                console.print(f"{i}. [{status}] {task}")

            choice = console.input("\nMark task(s) (e.g., '1 3 5') or [done]: ").strip()
            if choice.lower() == 'done':
                break

            try:
                # Mark multiple tasks at once
                selected_tasks = [int(num)-1 for num in choice.split()]
                for task_num in selected_tasks:
                    if 0 <= task_num < len(tasks):
                        completed.add(task_num)
                        progress.update(task_progress, advance=1)
                    else:
                        console.print(f"[red]Invalid task number: {task_num+1}[/red]")
            except ValueError:
                console.print("[red]Invalid input! Enter numbers separated by spaces.[/red]")

    completion_percentage = len(completed) / len(tasks) * 100
    return completion_percentage

def export_roadmap(roadmap, format_type="pdf"):
    """
    Saves roadmap as PDF/TXT
    """
    filename = f"study_plan_{len(roadmap)}.{format_type}"
    
    if format_type == "pdf":
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, roadmap)
        pdf.output(filename)
        
    elif format_type == "txt":
        with open(filename, 'w') as f:
            f.write(roadmap)
            
    console.print(f"[green]Saved as {filename}![green]")

def show_achievements(progress):
    """
    Displays achievements and points earned
    """
    points = int(progress * 10)  # 10 points per % completion

    console.print("\n[bold]Achievements Earned:[/]")
    if points >= 500:
        console.print("- 🏆 Master Learner: You've completed over 50% of your roadmap!")
    elif points >= 200:
        console.print("- 🧅 Dedicated Scholar: You're halfway there! Keep going!")
    elif points >= 100:
        console.print("- 🥉 Consistent Starter: Great start! Keep the momentum!")
    else:
        console.print("- 🌱 Beginner: You're just getting started!")

    console.print(f"\n[bold]Points Earned:[/] {points}")
    
    # Show next milestone
    if points < 100:
        console.print(f"[bold]Next Milestone:[/] Reach 100 points (complete {10 - points//10} more tasks)!")
    elif points < 200:
        console.print(f"[bold]Next Milestone:[/] Reach 200 points (complete {20 - points//10} more tasks)!")
    elif points < 500:
        console.print(f"[bold]Next Milestone:[/] Reach 500 points (complete {50 - points//10} more tasks)!")
    else:
        console.print("[bold]Next Milestone:[/] You've reached the top! 🎉")

def main():
    doc_content = ""
    roadmap = ""
    
    while True:
        console.print("\n[bold]Study Planner v2.0[/]")
        console.print("1. Generate Roadmap")
        console.print("2. Upload Documents (PDF/PPT)")
        console.print("3. Chat with Study Buddy")
        console.print("4. Track Progress")
        console.print("5. Export Roadmap")
        console.print("6. Exit")
        
        choice = console.input("\nChoose option (1-6): ")
        
        if choice == '1':
            user_input = console.input("\nEnter your goal (e.g., 'Learn ML in 10 days'): ")
            roadmap = generate_roadmap(user_input, doc_content)
            console.print(Markdown("\n**Your Roadmap:**\n" + roadmap))
            
        elif choice == '2':
            file_path = console.input("\nEnter document path: ").strip('"')
            doc_content = process_document(file_path)
            console.print(f"[green]Processed {len(doc_content)} characters![/green]")
            
        elif choice == '3':
            if roadmap or doc_content:
                study_chatbot(roadmap, doc_content)
            else:
                console.print("[red]Generate roadmap or upload documents first![/red]")
                
        elif choice == '4':
            if roadmap:
                progress = track_progress(roadmap)
                show_achievements(progress)
            else:
                console.print("[red]Generate roadmap first![/red]")
                
        elif choice == '5':
            if roadmap:
                fmt = console.input("Format (pdf/txt): ").lower()
                export_roadmap(roadmap, fmt)
            else:
                console.print("[red]Generate roadmap first![/red]")
                
        elif choice == '6':
            break
            
        else:
            console.print("[red]Invalid choice![/red]")
if __name__ == '__main__':
    main()